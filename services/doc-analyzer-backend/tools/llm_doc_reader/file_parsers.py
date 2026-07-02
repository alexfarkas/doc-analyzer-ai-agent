import csv
import json
import logging
from pathlib import Path

import chardet
import fitz  # PyMuPDF
import yaml
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from tools.llm_doc_reader.llm_doc_reader_config import LLMDocReaderConfig
from tools.llm_doc_reader.consts import CODE_LANG_MAP, RE_EMPTY_LINES, RE_WHITESPACE
from tools.llm_doc_reader.utils import (
    assemble_markdown,
    table_to_md,
    extract_office_images,
    run_ocr,
)

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────────────────────
#  Файловые парсеры
# ──────────────────────────────────────────────────────────────
class FileParser:
    def __init__(self, config: LLMDocReaderConfig):
        self.config = config

    def read_file(self, path: Path) -> str:
        ext = path.suffix.lower()
        handlers = {
            ".txt": self._parse_text,
            ".md": self._parse_text,
            ".csv": self._parse_csv,
            ".json": self._parse_structured,
            ".yaml": self._parse_structured,
            ".yml": self._parse_structured,
            ".html": self._parse_markup,
            ".xml": self._parse_markup,
            ".docx": self._parse_docx,
            ".xlsx": self._parse_xlsx,
            ".pptx": self._parse_pptx,
            ".pdf": self._parse_pdf,
            ".jpg": self._parse_image,
            ".jpeg": self._parse_image,
            ".png": self._parse_image,
            ".gif": self._parse_image,
            ".tiff": self._parse_image,
            ".bmp": self._parse_image,
            ".py": self._parse_code,
            ".js": self._parse_code,
            ".ts": self._parse_code,
            ".java": self._parse_code,
            ".kt": self._parse_code,
            ".scala": self._parse_code,
            ".cs": self._parse_code,
            ".cpp": self._parse_code,
            ".go": self._parse_code,
            ".php": self._parse_code,
            ".swift": self._parse_code,
            ".r": self._parse_code,
            ".pl": self._parse_code,
            ".sql": self._parse_code,
            ".sh": self._parse_code,
            ".zsh": self._parse_code,
            ".bash": self._parse_code,
        }
        # Code fallback для неизвестных расширений
        handler = handlers.get(ext) or (
            self._parse_code if ext in CODE_LANG_MAP else None
        )
        if not handler:
            logger.error(
                f"Unsupported extension: {ext}. Supported: {', '.join(handlers.keys())}"
            )

        blocks, images = handler(path)
        return assemble_markdown(blocks, images, path, self.config)

    def _parse_text(self, path: Path) -> tuple[list[str], list]:
        raw = path.read_bytes()
        enc = chardet.detect(raw)["encoding"] or "utf-8"
        text = raw.decode(enc, errors="replace")
        return [RE_EMPTY_LINES.sub("\n\n", RE_WHITESPACE.sub(" ", text).strip())], []

    def _parse_csv(self, path: Path) -> tuple[list[str], list]:
        with open(path, "r", encoding="utf-8-sig") as f:
            rows = [r for r in csv.reader(f) if any(c.strip() for c in r)]
        if not rows:
            return [""], []
        md = " | ".join(rows[0]) + "\n" + " | ".join(["---"] * len(rows[0]))
        md += "\n" + "\n".join(" | ".join(r) for r in rows[1:])
        return [f"CSV-таблица\n\n{md}"], []

    def _parse_structured(self, path: Path) -> tuple[list[str], list]:
        raw = path.read_text(encoding="utf-8-sig", errors="replace")
        data = json.loads(raw) if path.suffix == ".json" else yaml.safe_load(raw)
        return [
            json.dumps(data, indent=2, ensure_ascii=False)
            if path.suffix == ".json"
            else yaml.dump(data, default_flow_style=False, allow_unicode=True)
        ], []

    def _parse_markup(self, path: Path) -> tuple[list[str], list]:
        html = path.read_text(encoding="utf-8-sig", errors="replace")
        soup = BeautifulSoup(html, "lxml")
        for tag in soup(
            ["script", "style", "nav", "footer", "header", "iframe", "noscript"]
        ):
            tag.decompose()
        text = RE_EMPTY_LINES.sub(
            "\n\n", soup.get_text(separator="\n", strip=True)
        ).strip()
        return [text], []

    def _parse_code(self, path: Path) -> tuple[list[str], list]:
        raw = path.read_bytes()
        enc = chardet.detect(raw)["encoding"] or "utf-8"
        text = raw.decode(enc, errors="replace")
        lang = CODE_LANG_MAP.get(path.suffix.lower(), "plaintext")
        return [f"```{lang}\n{text.strip()}\n```"], []

    def _parse_docx(self, path: Path) -> tuple[list[str], list]:
        doc = Document(path)
        blocks, images = [], []
        for p in doc.paragraphs:
            if not p.text.strip():
                continue
            style = (p.style.name or "").lower()
            if "heading" in style:
                lvl = 1 if "1" in style else 2 if "2" in style else 3
                blocks.append(f"{'#' * lvl} {p.text.strip()}")
            elif "list" in style:
                blocks.append(f"- {p.text.strip()}")
            else:
                blocks.append(p.text.strip())
        for i, t in enumerate(doc.tables):
            if md := table_to_md(t):
                blocks.append(f"Таблица {i + 1}\n\n{md}")
        images = extract_office_images(path, "word/media/")
        return blocks, images

    def _parse_xlsx(self, path: Path) -> tuple[list[str], list]:
        wb = load_workbook(path, read_only=True, data_only=True)
        blocks = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows = [
                r
                for r in ws.iter_rows(values_only=True)
                if any(str(c).strip() for c in r if c is not None)
            ]
            if not rows:
                continue
            hdr = [str(h).strip() for h in rows[0]]
            hdr = [h if h else f"Col_{i}" for i, h in enumerate(hdr)]
            tbl = (
                f"### Лист: {sheet}\n\n"
                + " | ".join(hdr)
                + "\n"
                + " | ".join(["---"] * len(hdr))
            )
            for row in rows[1:]:
                tbl += "\n" + " | ".join(
                    str(c).replace("|", "\\|").strip()[:150] if c else "" for c in row
                )
            blocks.append(tbl)
        wb.close()
        return blocks, []

    def _parse_pptx(self, path: Path) -> tuple[list[str], list]:
        prs = Presentation(path)
        blocks, images = [], []
        for idx, slide in enumerate(prs.slides, 1):
            s_blocks = [f"# Слайд {idx}"]
            if slide.shapes.title and slide.shapes.title.text.strip():
                s_blocks[0] += f": {slide.shapes.title.text.strip()}"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if txt := p.text.strip():
                            s_blocks.append(txt)
                if shape.has_table:
                    if md := table_to_md(shape.table):
                        s_blocks.append(f"Таблица\n\n{md}")
                if hasattr(shape, "image") and shape.image:
                    images.append({"bytes": shape.image.blob, "idx": len(images) + 1})
            if slide.has_notes_slide:
                try:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        s_blocks.append(f"> Примечания: {notes}")
                except:
                    pass
            blocks.append("\n\n".join(b for b in s_blocks if b))
        return blocks, images

    def _parse_pdf(self, path: Path) -> tuple[list[str], list]:
        doc = fitz.open(path)
        blocks, images = [], []
        for i, page in enumerate(doc):
            text = page.get_text("text", sort=True).strip()
            page_blocks = [f"# Страница {i + 1}"]
            if text:
                page_blocks.append(text)
            else:
                page_blocks.append(
                    "[Нет текстового слоя. Будет выполнен локальный OCR...]"
                )
            for j, tbl in enumerate(page.find_tables().extract()):
                if md := table_to_md(tbl, is_pdf=True):
                    page_blocks.append(f"Таблица {j + 1}\n\n{md}")
            blocks.append("\n\n".join(page_blocks))
            if not text.strip() or len(text.strip()) < 50:
                pix = page.get_pixmap(dpi=200)
                images.append(
                    {
                        "bytes": pix.tobytes("png"),
                        "idx": len(images) + 1,
                        "type": "pdf_page_ocr",
                    }
                )
        if self.config.ocr_enabled and any("OCR" in b for b in blocks):
            blocks = []
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=200)
                blocks.append(
                    f"# Страница {i + 1} (OCR)\n\n{run_ocr(pix.tobytes('png'), self.config)}"
                )
        doc.close()
        return blocks, images

    def _parse_image(self, path: Path) -> tuple[list[str], list]:
        return [], [{"bytes": path.read_bytes(), "idx": 1}]

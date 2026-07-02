import base64
import csv
import json
import re
from pathlib import Path

import chardet
import fitz  # PyMuPDF
import yaml
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from components.preview_reader.preview_reader import CODE_LANGS, IMAGE_MIMES
from components.preview_reader.utils import normalize_rows, extract_zip_media


# ──────────────────────────────────────────────────────────────
#  Обработчики форматов
# ──────────────────────────────────────────────────────────────
def parse_text(path: Path) -> tuple[list[dict], dict]:
    raw = path.read_bytes()
    enc = chardet.detect(raw)["encoding"] or "utf-8"
    text = raw.decode(enc, errors="replace")
    text = re.sub(r"\r\n|\r", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return [{"type": "text", "content": text}], {
        "encoding": enc,
        "lines": text.count("\n"),
    }


def parse_markup(path: Path) -> tuple[list[dict], dict]:
    html = path.read_text(encoding="utf-8-sig", errors="replace")
    return parse_markup_text(html)


def parse_markup_text(text: str) -> tuple[list[dict], dict]:
    soup = BeautifulSoup(text, "lxml")
    for tag in soup(
        ["script", "style", "nav", "footer", "header", "iframe", "noscript", "form"]
    ):
        tag.decompose()
    title = soup.title.string.strip() if soup.title else ""
    text = re.sub(r"\n{3,}", "\n\n", soup.get_text(separator="\n", strip=True)).strip()
    return [{"type": "text", "content": text, "title": title}], {"title": title}


def parse_data(path: Path) -> tuple[list[dict], dict]:
    raw = path.read_text(encoding="utf-8-sig", errors="replace")
    lang = "yaml" if path.suffix.lower() in (".yaml", ".yml") else "json"
    data = yaml.safe_load(raw) if lang == "yaml" else json.loads(raw)
    formatted = (
        yaml.dump(data, default_flow_style=False, allow_unicode=True)
        if lang == "yaml"
        else json.dumps(data, indent=2, ensure_ascii=False)
    )
    return [{"type": "code", "language": lang, "content": formatted}], {"type": lang}


def parse_csv(path: Path) -> tuple[list[dict], dict]:
    with open(path, "r", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows = [r for r in reader if any(cell.strip() for cell in r)]
    if not rows:
        return [], {"rows": 0}
    return [{"type": "table", "headers": rows[0], "rows": rows[1:]}], {
        "rows": len(rows) - 1,
        "cols": len(rows[0]),
    }


def parse_code(path: Path) -> tuple[list[dict], dict]:
    raw = path.read_bytes()
    enc = chardet.detect(raw)["encoding"] or "utf-8"
    text = raw.decode(enc, errors="replace")
    lang = CODE_LANGS.get(path.suffix.lower(), "plaintext")
    return [{"type": "code", "language": lang, "content": text}], {"language": lang}


def parse_image(path: Path) -> tuple[list[dict], dict]:
    mime = IMAGE_MIMES.get(path.suffix.lower(), "application/octet-stream")
    b64 = base64.b64encode(path.read_bytes()).decode()
    src = f"data:{mime};base64,{b64}"
    return [{"type": "image", "src": src, "alt": path.name}], {
        "mime": mime,
        "size_kb": round(path.stat().st_size / 1024, 2),
    }


def parse_docx(path: Path) -> tuple[list[dict], dict]:
    doc = Document(path)
    blocks, imgs = [], []
    for p in doc.paragraphs:
        if not p.text.strip():
            continue
        style = (p.style.name or "").lower()
        if "heading" in style:
            lvl = 1 if "1" in style else 2 if "2" in style else 3
            blocks.append({"type": "heading", "level": lvl, "content": p.text.strip()})
        elif "list" in style:
            blocks.append({"type": "list_item", "content": p.text.strip()})
        else:
            blocks.append({"type": "text", "content": p.text.strip()})
    for i, t in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in r.cells] for r in t.rows]
        rows = normalize_rows(rows)
        if rows:
            blocks.append({"type": "table", "headers": rows[0], "rows": rows[1:]})

    # Извлечение картинок из ZIP-структуры .docx
    imgs = extract_zip_media(path, "word/media/")
    for img in imgs:
        blocks.append({"type": "image", "src": img["src"], "alt": img["name"]})
    return blocks, {
        "paragraphs": len(doc.paragraphs),
        "tables": len(doc.tables),
        "images": len(imgs),
    }


def parse_xlsx(path: Path) -> tuple[list[dict], dict]:
    wb = load_workbook(path, read_only=True, data_only=True)
    blocks = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = [
            r
            for r in ws.iter_rows(values_only=True)
            if any(str(c).strip() for c in r if c is not None)
        ]
        if not rows:
            continue
        rows = [[str(c).strip() if c is not None else "" for c in r] for r in rows]
        rows = normalize_rows(rows)
        blocks.append(
            {"type": "table", "sheet": sheet, "headers": rows[0], "rows": rows[1:]}
        )
    wb.close()
    return blocks, {"sheets": len(wb.sheetnames)}


def parse_pptx(path: Path) -> tuple[list[dict], dict]:
    prs = Presentation(path)
    blocks, img_count = [], 0
    for idx, slide in enumerate(prs.slides, 1):
        s_blocks = [{"type": "heading", "level": 1, "content": f"Слайд {idx}"}]
        if slide.shapes.title and slide.shapes.title.text.strip():
            s_blocks[-1]["content"] += f": {slide.shapes.title.text.strip()}"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if txt := p.text.strip():
                        s_blocks.append({"type": "text", "content": txt})
            if shape.has_table:
                rows = [
                    [cell.text.strip() for cell in r.cells] for r in shape.table.rows
                ]
                rows = normalize_rows(rows)
                if rows:
                    s_blocks.append(
                        {"type": "table", "headers": rows[0], "rows": rows[1:]}
                    )
            if hasattr(shape, "image") and shape.image:
                mime = f"image/{shape.image.ext.lower()}"
                b64 = base64.b64encode(shape.image.blob).decode()
                s_blocks.append(
                    {
                        "type": "image",
                        "src": f"data:{mime};base64,{b64}",
                        "alt": f"slide_{idx}_img_{img_count + 1}",
                    }
                )
                img_count += 1
        if slide.has_notes_slide:
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    s_blocks.append({"type": "note", "content": notes})
            except:
                pass
        blocks.append({"type": "slide_container", "blocks": s_blocks})
    return blocks, {"slides": len(prs.slides), "images": img_count}


def parse_pdf(path: Path) -> tuple[list[dict], dict]:
    doc = fitz.open(path)
    blocks = []
    for i, page in enumerate(doc):
        page_blocks = [{"type": "heading", "level": 1, "content": f"Страница {i + 1}"}]
        text = page.get_text("text", sort=True).strip()
        if text:
            for block in text.split("\n\n"):
                if b := block.strip():
                    page_blocks.append({"type": "text", "content": b})
        else:
            page_blocks.append(
                {"type": "text", "content": "[Текстовый слой отсутствует]"}
            )

        # Таблицы
        for j, tbl in enumerate(page.find_tables().extract()):
            rows = [[str(c).strip() if c else "" for c in r] for r in tbl]
            rows = normalize_rows(rows)
            if rows:
                page_blocks.append(
                    {"type": "table", "headers": rows[0], "rows": rows[1:]}
                )

        blocks.append({"type": "page_container", "blocks": page_blocks})
    doc.close()
    return blocks, {"pages": len(blocks)}
